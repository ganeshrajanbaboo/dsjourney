import pyodbc
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import date
from datetime import datetime, timedelta


class AddTable():
    def __init__(self, rows, cols, left, top, width, height, shapes):
        self.table = self.shapes.add_table(rows, cols, left, top, width, height).table
        
        
        

class SlideTemplate(AddTable):
    img_path = r'C:\Users\grajanbaboo\OneDrive - DXC Production\My Documents\WSR\Template\dxc.png'
    def __init__(self,prs):
        title_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(title_slide_layout)
        self.shapes = slide.shapes
        pic = self.shapes.add_picture(self.img_path, 1, 1)
        

class TaskTemplate(SlideTemplate):
    
    def cellformat(self, cell, cr, cg, cb, text, tf, fontname, fontsize, bold, fr, fg, fb):
        #format cell
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(cr, cg, cb)

        #format font
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        font = run.font
        font.name = fontname
        font.size = Pt(fontsize)
        font.bold = bold
        font.color.rgb = RGBColor(fr, fg, fb)

                
    def __init__(self,prs, inpdict, application, slidenum):
        SlideTemplate.__init__(self, prs)

        left = 0.5
        heading = ['Milestone', 'Start Date', 'Release Date', '% Complete', 'Status  ']
        totwidth = 0
        for k, v in inpdict.items():
            #v.insert(0,heading[0])
            #heading.pop(0)
            vx = v.copy()
            vx.insert(0,heading[0])
            maxstrlen = max([len(str(i)) for i in vx])
            if   k == 'Status':
                width = 2.2
            else:
                width = 3.0 if maxstrlen > 20 else maxstrlen * 0.12
            totwidth += width
            height = len(v) * 1
            #print('key = ', k)
            #print(v)
            AddTable.__init__(self, 1, 1, Inches(left), Inches(1.1), Inches(width), Inches(0.3), self.shapes)
            tf = self.table.cell(0, 0).text_frame
            self.cellformat(self.table.cell(0, 0), 0xE1, 0xE2, 0xE8, heading[0], tf, 'arial', 11, True, 0x12, 0x12, 0x12)
            heading.pop(0)
            AddTable.__init__(self, len(v), 1, Inches(left), Inches(1.7), Inches(width), Inches(height), self.shapes)
            left += width
            for r in range(len(v)):
                cell = self.table.cell(r, 0)
                tf = self.table.cell(r, 0).text_frame    
                if   k == 'Status' and v[r][:8] == 'Critical':
                    v[r] = v[r][8:]
                    cr, cg, cb = [0xED, 0x18, 0x18]
                elif k == 'Status' and v[r][:8] == 'Abnormal':
                    v[r] = v[r][8:]
                    cr, cg, cb = [0xE9, 0xED, 0x09]
                elif k == 'Status' and v[r][:11] == 'In Progress':
                    v[r] = v[r][11:]
                    cr, cg, cb = [0x67, 0xB5, 0x53]
                elif k == 'Status' and v[r][:9] == 'Completed':
                    v[r] = v[r][9:]
                    cr, cg, cb = [0x67, 0xB5, 0x53]
                else:
                    cr, cg, cb = [0xE1, 0xE2, 0xE8]
                    
                self.cellformat(cell, cr, cg, cb, str(v[r]), tf, 'calibri', 12, False, 0x12, 0x12, 0x12)
                
        #Add Heading            
        AddTable.__init__(self, 1, 1, Inches(0.5), Inches(0.8), Inches(totwidth), Inches(0.3), self.shapes)
        headtext = r'MF Tower Status update - ' + application + (r' - ' + str(slidenum) if slidenum > 1 else '' )
        tf = self.table.cell(0, 0).text_frame
        self.cellformat(self.table.cell(0, 0), 0x4C, 0x45, 0xA8, headtext, tf, 'arial', 12, True, 0xFF, 0xFF, 0xFF)

        #Add sub heading
        AddTable.__init__(self, 1, 1, Inches(0.5), Inches(1.4), Inches(totwidth), Inches(0.3), self.shapes)
        subheadtext = 'Development Activities'
        tf = self.table.cell(0, 0).text_frame
        self.cellformat(self.table.cell(0, 0), 0x8D, 0x98, 0xA1, subheadtext, tf, 'calibri', 11, True, 0x12, 0x12, 0x12)

        #Add Status column legend
        legendtop = height + 1.8
        AddTable.__init__(self, 1, 4, Inches(0.5), Inches(legendtop), Inches(5.0), Inches(0.3), self.shapes)
        tf = self.table.cell(0, 0).text_frame
        self.cellformat(self.table.cell(0, 0), 0x8D, 0x98, 0xA1, 'Status color legend = ', tf, 'arial', 11, True, 0x12, 0x12, 0x12)
        tf = self.table.cell(0, 1).text_frame
        self.cellformat(self.table.cell(0, 1), 0x67, 0xB5, 0x53, 'Normal', tf, 'calibri', 11, False, 0x12, 0x12, 0x12)
        tf = self.table.cell(0, 2).text_frame
        self.cellformat(self.table.cell(0, 2), 0xE9, 0xED, 0x09, 'Abnormal', tf, 'calibri', 11, False, 0x12, 0x12, 0x12)
        tf = self.table.cell(0, 3).text_frame
        self.cellformat(self.table.cell(0, 3), 0xED, 0x18, 0x18, 'Critical', tf, 'calibri', 11, False, 0x12, 0x12, 0x12)


        #self.shapes.title.text = 'Adding a table'
        #rows = cols = 2
        #left = top = Inches(2.0)
        #width = Inches(6.0)
        #height = Inches(0.8)
        #table = self.shapes.add_table(rows, cols, left, top, width, height).table

        #set column widths
        #table.columns[0].width = Inches[2.0]
        #table.columns[1].width = Inches[4.0]

        #write column headings
        #table.cell(0, 0).text = 'Foo'
        #table.cell(0, 1).text = 'Bar'

        #Write body cells
        #table.cell(1, 0).text = application
        #cell = table.cell(1, 0)
        #fill = cell.fill
        #fill.solid()
        #fill.fore_color.rgb = RGBColor(0x3B, 0xCA, 0x79)
        #table.cell(1, 1).text = 'Qux'


class BatchTemplate(TaskTemplate):
    
    def __init__(self,prs, inpdict, application, slidenum):
        SlideTemplate.__init__(self, prs)

        left = 1.2
        lefthead = left
        heading = ['Exec Date', 'Region      ', 'Nature', 'Status', 'Run By   ']
        totwidth = 0
        for k, v in inpdict.items():
            #v.insert(0,heading[0])
            #heading.pop(0)
            vx = v.copy()
            vx.insert(0,heading[0])
            maxstrlen = max([len(str(i)) for i in vx])
            width = 2.5 if maxstrlen > 25 else (2.5/25)*maxstrlen
            totwidth += width
            height = len(v) * 0.7
            #print('key = ', k)
            #print(v)
            AddTable.__init__(self, 1, 1, Inches(left), Inches(1.1), Inches(width), Inches(0.3), self.shapes)
            tf = self.table.cell(0, 0).text_frame
            self.cellformat(self.table.cell(0, 0), 0xE1, 0xE2, 0xE8, heading[0], tf, 'arial', 11, True, 0x12, 0x12, 0x12)
            heading.pop(0)
            AddTable.__init__(self, len(v), 1, Inches(left), Inches(1.4), Inches(width), Inches(height), self.shapes)
            left += width
            for r in range(len(v)):
                cell = self.table.cell(r, 0)
                tf = self.table.cell(r, 0).text_frame
                cr, cg, cb = [0xE1, 0xE2, 0xE8]
                self.cellformat(cell, cr, cg, cb, str(v[r]), tf, 'calibri', 11, False, 0x12, 0x12, 0x12)
                
        #Add Heading            
        AddTable.__init__(self, 1, 1, Inches(lefthead), Inches(0.8), Inches(totwidth), Inches(0.3), self.shapes)
        headtext = r'Batch Cycle Report - ' + application + (r' - ' + str(slidenum) if slidenum > 1 else '' )
        tf = self.table.cell(0, 0).text_frame
        self.cellformat(self.table.cell(0, 0), 0x4C, 0x45, 0xA8, headtext, tf, 'arial', 12, True, 0xFF, 0xFF, 0xFF)



class UtilTemplate(TaskTemplate):
    def __init__(self,prs, inpdict, slidenum):
        SlideTemplate.__init__(self, prs)

        left = 0.5
        heading = ['Person', 'Utilization hours (for reported week)', 'Holiday/Leave', 'Utilization hours forecast (for coming week)', 'Remarks']
        totwidth = 0
        for k, v in inpdict.items():
            vx = v.copy()
            vx.insert(0,heading[0])
            
            maxstrlen = max([len(str(i)) for i in vx])
            width = 2.0 if maxstrlen > 15 else (2.0/15)*maxstrlen
            totwidth += width
            height = len(v) * 0.5
            AddTable.__init__(self, 1, 1, Inches(left), Inches(1.1), Inches(width), Inches(0.6), self.shapes)
            tf = self.table.cell(0, 0).text_frame
            TaskTemplate.cellformat(self, self.table.cell(0, 0), 0x8D, 0x98, 0xA1, heading[0], tf, 'arial', 11, True, 0x12, 0x12, 0x12)
            heading.pop(0)
            AddTable.__init__(self, len(v), 1, Inches(left), Inches(1.7), Inches(width), Inches(height), self.shapes)
            left += width
            for r in range(len(v)):
                cell = self.table.cell(r, 0)
                tf = self.table.cell(r, 0).text_frame
                cr, cg, cb = [0xE1, 0xE2, 0xE8]
                TaskTemplate.cellformat(self, cell, cr, cg, cb, str(v[r]), tf, 'calibri', 11, False, 0x12, 0x12, 0x12)
                
        #Add Heading            
        AddTable.__init__(self, 1, 1, Inches(0.5), Inches(0.8), Inches(totwidth), Inches(0.3), self.shapes)
        headtext = r'Utilization' + (r' - ' + str(slidenum) if slidenum > 1 else '' )
        tf = self.table.cell(0, 0).text_frame
        TaskTemplate.cellformat(self, self.table.cell(0, 0), 0x4C, 0x45, 0xA8, headtext, tf, 'arial', 12, True, 0xFF, 0xFF, 0xFF)


class MiscTemplate(TaskTemplate):
    def __init__(self,prs, inpdict, slidenum):
        SlideTemplate.__init__(self, prs)
        top =  1.3
        for k, v in inpdict.items():
            prevtop = top
            keyheight = 0
            for r in range(len(v)):
                AddTable.__init__(self, 1, 1, Inches(3.0), Inches(top), Inches(6.0), Inches(0.6), self.shapes)
                cell = self.table.cell(0, 0)
                tf = self.table.cell(0, 0).text_frame
                cr, cg, cb = [0xE1, 0xE2, 0xE8]
                TaskTemplate.cellformat(self, cell, cr, cg, cb, str(v[r]), tf, 'calibri', 11, False, 0x12, 0x12, 0x12)
                keyheight += 0.6
                top += 0.6
            top += 0.1
            #height += top
            AddTable.__init__(self, 1, 1, Inches(0.5), Inches(prevtop), Inches(2.5), Inches(keyheight), self.shapes)
            cell = self.table.cell(0, 0)
            tf = self.table.cell(0, 0).text_frame
            cr, cg, cb = [0x8D, 0x98, 0xA1]
            TaskTemplate.cellformat(self, cell, cr, cg, cb, str(k), tf, 'arial', 11, True, 0x12, 0x12, 0x12)
            
        #Add Heading            
        AddTable.__init__(self, 1, 1, Inches(0.5), Inches(0.8), Inches(8.5), Inches(0.3), self.shapes)
        headtext = r'Miscellaneous' + (r' - ' + str(slidenum) if slidenum > 1 else '' )
        tf = self.table.cell(0, 0).text_frame
        TaskTemplate.cellformat(self, self.table.cell(0, 0), 0x4C, 0x45, 0xA8, headtext, tf, 'arial', 12, True, 0xFF, 0xFF, 0xFF)
            
            

################################### E N D  O F  C L A S S  D E F I N I T I O N S ###########################################

#        
#B E G I N   O F  F U N C T I O N S  C A L L E D  F R O M  M A I N
#
def readtable(crsr, tablename, num_cols, wherepredicate, params):
    #initialize the return dictionary
    outdict = {}

    #Create the select query and execute it
    selectquery = r'select * from ' + tablename + wherepredicate
    #print(selectquery)
    crsr.execute(selectquery, params)

    #Initialize a column names list
    columnnames = []

    #Initialize a set of column values lists
    columnvalues = [[] for i in range(num_cols)]

    #Fetch first row
    firstrow = crsr.fetchone()
    if firstrow == None:
        return outdict
    for i in range (0, num_cols):
        #Build a list of columns.
        #Elements of this list will serve as keys to a dictionary of data
        columnnames.append(firstrow.cursor_description[i][0])
        #Get the values for the first row fetched in the first of the set of column lists
        columnvalues[i].append(firstrow[i])

    #Fetch all the remaining rows
    for row in crsr.fetchall():
        for i in range (0, num_cols):
            columnvalues[i].append(row[i])

    #Build a dictionary of column names key and column values data and return it
    outdict = {}
    #rename duplicates keeping the CPSIN_Utilization query in mind
    if len(columnnames) != len(set(columnnames)):
        columnnamesuniq = []
        for i in columnnames:
            #remove duplicates if any
            [columnnamesuniq.append(i) for i in columnnames if not i in columnnamesuniq]
        columnnames.clear()
        for x in columnnamesuniq:
            columnnames.append('A.' + x)
        for y in columnnamesuniq:
            columnnames.append('B.' + y)

    #print('COLUMNS')
    #print(columnnames)
    #print('VALUES')
    #print(columnvalues)
    
    for i in enumerate(columnnames):
        #print(columnnames[i[0]])
        outdict.update({columnnames[i[0]]: columnvalues[i[0]]})
    return outdict             
    

def combinedict(dict1, dict2):
    if dict1 == {}:
        return dict2
    else:
        dict3 = dict1
    for k, v in dict1.items():
        #print(k, v)
        if k in dict2:
            dict3[k] = dict1[k] + dict2[k]
    return dict3

def listsort(list1, list2):
    list3, list4 = (list(t) for t in zip(*sorted(zip(list1, list2))))
    return list4

def defaultval(listx):
    if listx.count(None) == len(listx):
        return listx
    elif listx.count(None) == 0:
        return listx
    else:
        if any(isinstance(x, int) for x in listx):
            return [0 if x == None else x for x in listx]
        elif any(isinstance(x, datetime) for x in listx):
            return [datetime.strptime('0001-01-01', '%Y-%m-%d') if x == None else x for x in listx]
        elif any(isinstance(x, str) for x in listx):
            return ['' if x == None else x for x in listx]

#extract the data
def dataextract(usedate):

    #Work Date
    #usedate = '2019-06-10'
    #format date range
    daterangestart = datetime.strptime(usedate, '%Y-%m-%d').date() - timedelta(days = 7)
    daterangeend = datetime.strptime(usedate, '%Y-%m-%d').date() - timedelta(days = 3)
    #create an empty tuple of params.
    params = ()

    #Initialize dictionaries that will be used to hold data as key-value pairs
    CPSIN_Tasks_dict = {}
    CPSIN_Batch_dict = {}
    CPSIN_Misc_dict  = {}
    CPSIN_Util_dict  = {}
    
    
    
    #Establish connection to MS Access Database
    conn_str = (
        r'DRIVER={Microsoft Access Driver (*.mdb, *.accdb)};'
        r'DBQ=C:\Users\grajanbaboo\OneDrive - DXC Production\My Documents\WSR\database\wsr1.accdb;'
        )
    #Start an ODBC connection between the program and the database using the connection string
    cnxn = pyodbc.connect(conn_str)
    #Open a cursor to interact with the database
    crsr = cnxn.cursor()
    #Fetch all table names in the database
    tablenames = [x for x in crsr.tables(tableType='TABLE')]
    #Browse through the names of the tables in the database
    for table_info in tablenames:
        #Check if the table names start with 'CPSIN_Tasks_' and read it
        if 'CPSIN_Tasks_' in table_info.table_name:
            #give the number of columns to read.
            num_cols = 8
            #Build the parameter tuple
            params = (daterangestart, daterangeend)
            #Build where predicate for 'CPSIN_Tasks_' table
            whereCPSIN_Tasks_ = r' where End_Date >=  ? AND End_Date <= ? OR Completion_Percentage < 100'
            #Read the table using the table name, number of columns and the where predicate and add to the corresponding dictionary
            CPSIN_Tasks_dict = combinedict(CPSIN_Tasks_dict, readtable(crsr, table_info.table_name, num_cols, whereCPSIN_Tasks_,  params))
            #print(CPSIN_Tasks_dict)
            
            
        #Read the CPSIN_Batch table
        elif table_info.table_name == 'CPSIN_Batch':
            #give the number of columns to read.
            num_cols = 7
            #Build the parameter tuple
            params = (daterangestart, daterangeend)
            #Build where predicate for 'CPSIN_Batch' table
            whereCPSIN_Batch = r' where Exec_Date >=  ? AND Exec_Date <= ? order by Exec_Date'
            #Read the table using the table name, number of columns and the where predicate and add to the corresponding dictionary
            CPSIN_Batch_dict = combinedict(CPSIN_Batch_dict, readtable(crsr, table_info.table_name, num_cols, whereCPSIN_Batch, params))
            
            
        #Read the CPSIN_Miscellaneous table
        elif table_info.table_name == 'CPSIN_Miscellaneous':
            #give the number of columns to read.
            num_cols = 4
            #initialize params tuple
            params = ()
            #Build where predicate for 'CPSIN_Miscellaneous' table
            whereCPSIN_Misc = " where Status = 'Open/Ongoing' order by Type"
            #Read the table using the table name, number of columns and the where predicate and add to the corresponding dictionary
            CPSIN_Misc_dict = combinedict(CPSIN_Misc_dict, readtable(crsr, table_info.table_name, num_cols, whereCPSIN_Misc,  params))
            
            
        #Read the CPSIN_Utilization table
        elif table_info.table_name == 'CPSIN_Utilization':
            #give the number of columns to read.
            num_cols = 7
            #format start week
            reportingstartweek   = datetime.strptime(usedate, '%Y-%m-%d').date()
            #reportingstartweek   = datetime.strptime(usedate, '%Y-%m-%d').date() - timedelta(days = 7)
            #forecastingstartweek = datetime.strptime(usedate, '%Y-%m-%d').date() 
            #build the params tuple
            params = reportingstartweek
            #<depricated>Create a FROM clause of self inner join
            #modtable = table_info.table_name  + r' AS A INNER JOIN ' +  table_info.table_name + r' AS B ON A.[Person] = B.[Person]'
            #Build where predicate of CPSIN_Utilization table
            whereCPSIN_Util = ' where Start_Week = ?'
            #Read the table using the table name, number of columns and the where predicate and add to the corresponding dictionary
            CPSIN_Util_dict = combinedict(CPSIN_Util_dict, readtable(crsr, table_info.table_name, num_cols, whereCPSIN_Util,  params))

    return CPSIN_Tasks_dict, CPSIN_Batch_dict, CPSIN_Misc_dict, CPSIN_Util_dict
        
           
def CPSINTasksSlides(prs, CPSIN_Tasks_dict, ns):

    #Merge the Status and the remarks columns into the Status column
    mergelist = [i + (str(j) if j != None else '')  for i, j in zip(CPSIN_Tasks_dict['Status'], CPSIN_Tasks_dict['Remarks'])]
    CPSIN_Tasks_dict['Status'].clear()
    CPSIN_Tasks_dict['Status'] = mergelist

    #Drop the Remarks column
    del CPSIN_Tasks_dict['Remarks']
    
    #Sort the Application list of the dictionary CPSIN_Tasks_dict and sort the remaining lists in that order
    worklist = CPSIN_Tasks_dict['Application']
    worklist1 = []
    for i, k in enumerate(worklist):
        k = k + str(i)
        worklist1.append(k)
    
    for k, v in CPSIN_Tasks_dict.items():
        v = defaultval(v)
        v = listsort(worklist1, v)
        CPSIN_Tasks_dict[k] = v

    
    #get the unique individual applications provided by the key 'Application' in a list
    appl = set(CPSIN_Tasks_dict['Application'])
    appl = sorted(list(appl))
        
    #capture the number of times each of those applications occur in a list
    applcount = []
    for elem in appl:
        applcount.append(CPSIN_Tasks_dict['Application'].count(elem))

    #Extract data for the slides
    startpos = 0
    for i, j in zip(range(len(applcount)), applcount):
        #startpos += (0 if i == 0 else applcount[i - 1])
        #number of slides required for each application and the data that goes in each slide
        for m in range(int(j/ns) if j % ns == 0 else int(j/ns) + 1):
            #startpos += m
            #number of elements to be included in the current instantiation of TaskTemplate            
            n = (ns if (j - m*ns) / ns >= 1.0 else (j - m*ns) % ns)
            #print("ns", "j", "m", "n", "sp")
            #print(ns, j, m, n, startpos)
            #slice CPSIN_Tasks_dict for a current instantiation of TaskTemplate
            inpdict = {}
            for k, v in CPSIN_Tasks_dict.items():
                if k in ('Task_ID', 'Application'):
                    continue
                else:
                    vwork = v[startpos : startpos + n]
                    if k in ('Start_Date', 'End_Date'):
                        vwork = [i if i is not None else datetime.strptime('0001-01-01', '%Y-%m-%d') for i in vwork]
                        vwork = [i.date() if str(i.date()) > '0001-01-01' else '' for i in vwork]
                    inpdict.update({k : vwork})
            startpos += n 

            #Create the slide
            out = TaskTemplate(prs, inpdict, appl[i], m+1)

def CPSINBatchSlides(prs, CPSIN_Batch_dict, ns):
    #Sort the Application list of the dictionary CPSIN_Batch_dict and sort the remaining lists in that order
    worklist = [i + str(j) for i, j in zip(CPSIN_Batch_dict['Application'], CPSIN_Batch_dict['Exec_Date'])] 
    worklist1 = []
    for i, k in enumerate(worklist):
        k = k + str(i)
        worklist1.append(k)
        
    for k, v in CPSIN_Batch_dict.items():
        v = defaultval(v)
        v = listsort(worklist1, v)
        CPSIN_Batch_dict[k] = v

    
    #get the unique individual applications provided by the key 'Application' in a list
    appl = set(CPSIN_Batch_dict['Application'])
    appl = sorted(list(appl))
        
    #capture the number of times each of those applications occur in a list
    applcount = []
    for elem in appl:
        applcount.append(CPSIN_Batch_dict['Application'].count(elem))


    #Extract data for the slides
    startpos = 0
    for i, j in zip(range(len(applcount)), applcount):
        #startpos += (0 if i == 0 else applcount[i - 1])
        #number of slides required for each application and the data that goes in each slide
        for m in range(int(j/ns) if j % ns == 0 else int(j/ns) + 1):
            #startpos += m
            #number of elements to be included in the current instantiation of TaskTemplate            
            n = (ns if (j - m*ns) / ns >= 1.0 else (j - m*ns) % ns)
            #print("ns", "j", "m", "n", "sp")
            #print(ns, j, m, n, startpos)
            #slice CPSIN_Tasks_dict for a current instantiation of TaskTemplate
            inpdict = {}
            for k, v in CPSIN_Batch_dict.items():
                if k in ('ID', 'Application'):
                    continue
                else:
                    vwork = v[startpos : startpos + n]
                    if k == 'Exec_Date':
                        vwork = [i if i is not None else datetime.strptime('0001-01-01', '%Y-%m-%d') for i in vwork]
                        vwork = [i.date() if str(i.date()) > '0001-01-01' else '' for i in vwork]
                    inpdict.update({k : vwork})
            startpos += n 

            #Create the slide
            out = BatchTemplate(prs, inpdict, appl[i], m+1)
        


def CPSINUtilSlides(prs, CPSIN_Util_dict, ns):
    #------------------------------------------Drop unwanted keys of CPSIN_Util_dict----------------------------------------------
    #print(CPSIN_Util_dict)
    #for key in ['A.ID', 'A.Start_Week', 'A.Encoded Absolute URL', 'A.Item Type', 'A.Path', 'A.URL Path',
    #                   'A.Workflow Instance ID', 'A.File Type', 'B.ID', 'B.Start_Week', 'B.Person', 'B.Holiday/Leave',
    #                   'B.Encoded Absolute URL', 'B.Item Type', 'B.Path', 'B.URL Path', 'B.Workflow Instance ID', 'B.File Type']:
    #   del CPSIN_Util_dict[key]
    #Combine Reporting week's and forecasting week's remarks.
    #UtilListB = [i for i in CPSIN_Util_dict['B.Remarks']]
    #CPSIN_Util_dict['B.Remarks'].clear()
    #for i, j in zip(UtilListB, CPSIN_Util_dict['A.Remarks']) :
    #   addstr = ''
    #   if j != None:
    #       addstr = ' reported: ' + j
    #   if i != None:
    #       addstr += ', forecasted: ' + i
    #   CPSIN_Util_dict['B.Remarks'].append(addstr)
        
    #Drop the key A.Remarks of CPSIN_Util_dict
    #del CPSIN_Util_dict['A.Remarks']

    #-------------------------------Slide creation logic-------------------------------
    #Drop the key ID of CPSIN_Util_dict
    del CPSIN_Util_dict['ID']
    del CPSIN_Util_dict['Start_Week']

    #number of people
    nop = len(CPSIN_Util_dict['Person'])

    #Create the slides
    #
    startpos = 0

    #number of slides required and the data that goes in each slide
    for m in range(int(nop/ns) if nop % ns == 0 else int(nop/ns) + 1):
        #number of elements to be included in the current instantiation of UtilTemplate
        n = (ns if nop % ns == 0 else nop % ns)

        #slice CPSIN_Util_dict for a current instantiation of UtilTemplate
        inpdict = {}
        startpos += m*ns
        for k, v in CPSIN_Util_dict.items():
            vwork = v[startpos : startpos + n]
            if k == "Remarks":
                vwork = [i if i is not None else "" for i in vwork]
            inpdict.update({k : vwork})

        #Create the slide
        out = UtilTemplate(prs, inpdict, m+1)
            

def CPSINMiscSlides(prs, CPSIN_Misc_dict, ns):      
    typecountdict = {}
    #Sort the types key in the order Issues, Appreciations, Training, Others

    misctypes = [t for t in CPSIN_Misc_dict['Type']]
    misctypes1 = []
    #Add sequence
    for t in misctypes:
        if t == 'Issues':
            t = str(1) + t
            
        elif t == 'Appreciations':
            t = str(2) + t
            
        elif t == 'Training':
            t = str(3) + t
            
        elif t == 'Others':
            t = str(4) + t
            
        else:
            t = str(5) + t
            
        misctypes1.append(t)

    misctypes2 = []
    for i, k in enumerate(misctypes1):
        k = k + str(i+100)
        misctypes2.append(k)

    #Sort each value list in the input dictionary in the order of misctypes2
    for k, v in CPSIN_Misc_dict.items():
        v = listsort(misctypes2, v)
        CPSIN_Misc_dict[k] = v

    #Remove last 3 characters (3 character sequence number) from misctypes2, which is a number
    misctypes3 = []
    for t in misctypes2:
        misctypes3.append(t[:-3])

    
    for i in sorted(set(misctypes3)):
        counti = CPSIN_Misc_dict['Type'].count(i[1:])
        typecountdict.update({(i[1:]):counti})

    typedatadict = {}
    startpos = 0
    endpos = 0
    
    for k, v in typecountdict.items():
        endpos = startpos + v
        vwork = CPSIN_Misc_dict['Remarks'][startpos : endpos]
        typedatadict.update({k : vwork})
        startpos = startpos + v

        
    #total number of entries
    ntot = 0
    for v in typecountdict.values():
        ntot +=  v 
    
    #Number of slides required
    nslides = int(ntot/ns) if ntot % ns == 0 else int(ntot/ns) + 1

    #Create a list for the data that goes in each slide
    slidelist = []
    #for i in range(nslides):
    inpdict = {}
    itemcount = 0
    vinpdata = []


    #converting typedatadict to a list of key value pairs.
    kp = []
    for k, v in typedatadict.items():
        for i in v:
            kp.append((k, i))

    
    startpos = 0
    for m in range(nslides):
        startpos = m*ns
        endpos = startpos + (ns if m < nslides - 1 else ntot - m*ns)
        sl = []
        sl.extend(kp[startpos : endpos])
        
        inpdict = {}
        k = sl[0][0]
        v = []
        for i in sl:
            if k != i[0]:
                vc = v.copy()
                inpdict.update({k : vc})
                v.clear()
                k = i[0]
            v.append(i[1])
        vc = v.copy()
        inpdict.update({k : vc})

        out = MiscTemplate(prs, inpdict, m+1)


    
        
def main():
    
    #Instantiate a presentation class object
    prs = Presentation()
    
    #Extract the data
    CPSIN_Tasks_dict, CPSIN_Batch_dict, CPSIN_Misc_dict, CPSIN_Util_dict = dataextract('2020-02-10')

    
    #max number of data rows in a Tasks slide
    ns = 5
    #Create the slides for Tasks
    CPSINTasksSlides(prs, CPSIN_Tasks_dict, ns)


    #max number of data rows in a Utilization slide
    ns = 10
    #Create the slides for Utilization
    CPSINUtilSlides(prs, CPSIN_Util_dict, ns)

    
    #max number of data rows in a Miscellaneous items slide
    ns = 10
    #Create the slides for Miscellaneous items
    CPSINMiscSlides(prs, CPSIN_Misc_dict, ns)

    
    #max number of data rows in a Batch execution slide
    ns = 8
    #Create the slides for Batch execution
    CPSINBatchSlides(prs, CPSIN_Batch_dict, ns)


    #Save the presentation
    prs.save(r'C:\Users\grajanbaboo\OneDrive - DXC Production\My Documents\WSR\Template\b.pptx')
    
    
    
    

                 
if __name__ == "__main__":
    main()

    

    




